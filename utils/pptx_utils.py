from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import json

add_border_label_function = r'''
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def pt_to_emu(points: float) -> int:
    return int(points * 12700)

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def add_border_and_labels(
    prs,
    border_color=RGBColor(255, 0, 0),   # Red border for shapes
    border_width=Pt(2),                # 2-point border width
    label_outline_color=RGBColor(0, 0, 255),  # Blue outline for label circle
    label_text_color=RGBColor(0, 0, 255),     # Blue text color
    label_diameter_pt=40                       # Diameter of the label circle in points
):
    """
    Iterates over all slides and shapes in the Presentation 'prs', applies a 
    red border to each shape, and places a transparent (no fill), blue-outlined 
    circular label with a blue number in the center of each shape. Labels start 
    from 0 and increment for every shape that gets a border.

    Args:
        prs: The Presentation object to modify.
        border_color: RGBColor for the shape border color (default: red).
        border_width: The width of the shape border (Pt).
        label_outline_color: The outline color for the label circle (default: blue).
        label_text_color: The color of the label text (default: blue).
        label_diameter_pt: The diameter of the label circle, in points (default: 40).
    """
    label_diameter_emu = pt_to_emu(label_diameter_pt)  # convert diameter (points) to EMUs
    label_counter = 0  # Start labeling at 0
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            # Skip shapes that are labels themselves
            if shape.name.startswith("Label_"):
                continue

            try:
                # --- 1) Add red border to the shape (if supported) ---
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                # --- 2) Calculate center for the label circle ---
                label_left = shape.left + (shape.width // 2) - (label_diameter_emu // 2)
                label_top  = shape.top  + (shape.height // 2) - (label_diameter_emu // 2)

                # --- 3) Create label circle (an OVAL) in the center of the shape ---
                label_shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    label_left,
                    label_top,
                    label_diameter_emu,
                    label_diameter_emu
                )
                label_shape.name = f"Label_{label_counter}"  # so we can skip it later

                # **Make the circle completely transparent** (no fill at all)
                label_shape.fill.background()

                # **Give it a blue outline**
                label_shape.line.fill.solid()
                label_shape.line.fill.fore_color.rgb = label_outline_color
                label_shape.line.width = Pt(3)

                # --- 4) Add the label number (centered, blue text) ---
                tf = label_shape.text_frame
                tf.text = str(label_counter)
                paragraph = tf.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER

                run = paragraph.runs[0]
                font = run.font
                font.size = Pt(40)      # Larger font
                font.bold = True
                font.name = "Arial"
                font._element.get_or_change_to_solidFill()
                font.fill.fore_color.rgb = label_text_color
                # Record properties from the original shape and label text.
                labeled_elements[label_counter] = {
                    'left': f'{emu_to_inches(shape.left)} Inches',
                    'top': f'{emu_to_inches(shape.top)} Inches',
                    'width': f'{emu_to_inches(shape.width)} Inches',
                    'height': f'{emu_to_inches(shape.height)} Inches',
                    'font_size': f'{shape.text_frame.font.size} PT' if hasattr(shape, 'text_frame') else None,
                }

                # --- 5) Increment label counter (so every shape has a unique label) ---
                label_counter += 1

            except Exception as e:
                # If the shape doesn't support borders or text, skip gracefully
                print(f"Could not add border/label to shape (type={shape.shape_type}): {e}")

    return labeled_elements
'''

add_border_function = r'''
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def add_border(
    prs,
    border_color=RGBColor(255, 0, 0),   # Red border for shapes
    border_width=Pt(2),                # 2-point border width
):
    """
    Iterates over all slides and shapes in the Presentation 'prs', applies a 
    red border to each shape, and places a transparent (no fill).

    Args:
        prs: The Presentation object to modify.
        border_color: RGBColor for the shape border color (default: red).
        border_width: The width of the shape border (Pt).
    """
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                # --- 1) Add red border to the shape (if supported) ---
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                if hasattr(shape, 'name'):
                    labeled_elements[shape.name] = {
                        'left': f'{emu_to_inches(shape.left)} Inches',
                        'top': f'{emu_to_inches(shape.top)} Inches',
                        'width': f'{emu_to_inches(shape.width)} Inches',
                        'height': f'{emu_to_inches(shape.height)} Inches',
                    }

            except Exception as e:
                # If the shape doesn't support borders or text, skip gracefully
                print(f"Could not add border to shape (type={shape.shape_type}): {e}")
    
    return labeled_elements
'''

create_id_map_function = r'''
def create_element_id_map(presentation):
    """
    Given a python-pptx Presentation object, this function creates
    and returns a dictionary mapping each element's (shape's) unique id
    to a sequential integer starting from 0.
    
    Parameters:
        presentation (Presentation): A python-pptx Presentation object.
        
    Returns:
        dict: A dictionary with keys as element IDs (integers) and values as sequential integers.
    """
    element_id_map = {}
    counter = 0

    # Iterate over each slide in the presentation
    for slide in presentation.slides:
        # Iterate over each shape (element) on the slide
        for shape in slide.shapes:
            if hasattr(shape, "name"):
                element_id_map[counter] = shape.name
                counter += 1

    return element_id_map
'''

save_helper_info_border_label = r'''
location_info = add_border_and_labels(poster, label_diameter_pt=80)
id_map = create_element_id_map(poster)
import json

with open('{}_element_id_map.json', 'w') as f:
    json.dump(id_map, f)

with open('{}_location_info.json', 'w') as f:
    json.dump(location_info, f)

poster.save("{}_bordered.pptx")
'''

save_helper_info_border = r'''
location_info = add_border(poster)
import json

with open('{}_location_info.json', 'w') as f:
    json.dump(location_info, f)

poster.save("{}_bordered.pptx")
'''

utils_functions = r'''

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import pptx
import json

from pptx.enum.text import MSO_AUTO_SIZE

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def _px_to_pt(px):
    """
    Approximate conversion from pixels to points.
    A common assumption is 1px ~ 0.75pt.
    Adjust as needed for your environment.
    """
    return px * 0.75

def _parse_font_size(font_size):
    """
    Internal helper to convert a numeric font size (e.g., 12) 
    to a python-pptx Pt object. If it's already a Pt, return as-is.
    """
    if font_size is None:
        return None
    if isinstance(font_size, (int, float)):
        return Pt(font_size)
    return font_size  # Assume user provided a Pt object already

def _parse_alignment(alignment):
    """
    Internal helper to convert a string alignment (e.g., "left", "center") 
    to the corresponding PP_ALIGN constant. 
    Default to PP_ALIGN.LEFT if unrecognized or None.
    """
    if not isinstance(alignment, str):
        # If user passed None or something else, default to PP_ALIGN.LEFT
        return PP_ALIGN.LEFT

    alignment = alignment.lower().strip()
    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return alignment_map.get(alignment, PP_ALIGN.LEFT)

def create_poster(width_inch=48, height_inch=36):
    """
    Create a new Presentation object, set its slide size (e.g., 48x36 inches).
    
    :param width_inch: Float or int specifying width in inches (default 48).
    :param height_inch: Float or int specifying height in inches (default 36).
    :return: A python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = Inches(width_inch)
    prs.slide_height = Inches(height_inch)
    return prs

def add_blank_slide(prs):
    """
    Add a blank slide to the Presentation (layout index 6 is typically blank).
    
    :param prs: The Presentation object to add a slide to.
    :return: The newly added slide object.
    """
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def shape_fill_color(shape, fill_color):
    """
    Set the fill color of a shape to the specified RGB color.

    :param shape: The shape object to modify.
    :param fill_color: A tuple (r, g, b) for the fill color.
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)


def add_textbox(
    slide, 
    name, 
    left_inch, 
    top_inch, 
    width_inch, 
    height_inch, 
    text="", 
    word_wrap=True,
    font_size=40,
    bold=False,
    italic=False,
    alignment="left",
    fill_color=None,
    font_name="Arial"
):
    """
    Create a textbox shape on the given slide, optionally fill its background with
    a color if fill_color is specified as (r, g, b).
    
    :param slide: Slide object to place the textbox on.
    :param name: Name for the shape (shape.name).
    :param left_inch: Left coordinate (in inches).
    :param top_inch: Top coordinate (in inches).
    :param width_inch: Width (in inches).
    :param height_inch: Height (in inches).
    :param text: Text to display in the textbox.
    :param word_wrap: If True, wrap text in the textbox.
    :param font_size: Numeric font size (e.g. 40).
    :param bold: Boolean to set run.font.bold.
    :param italic: Boolean to set run.font.italic.
    :param alignment: String alignment: "left", "center", "right", or "justify".
    :param fill_color: (r, g, b) tuple for solid fill background color, or None to skip.
    :param font_name: String font name (e.g., "Arial").
    :return: The newly created textbox shape.
    """
    shape = slide.shapes.add_textbox(
        Inches(left_inch), Inches(top_inch),
        Inches(width_inch), Inches(height_inch)
    )
    
    shape.name = name
    
    # If a fill color is specified, apply a solid fill
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        # Otherwise, set "no fill" if you want it transparent
        shape.fill.background()

    text_frame = shape.text_frame
    # Turn off auto-size to ensure stable font size, etc.
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.word_wrap = word_wrap
    
    # Clear any default paragraphs
    text_frame.clear()
    
    # Add a new paragraph
    p = text_frame.add_paragraph()
    # Instead of setting p.text, explicitly create a Run
    run = p.add_run()
    run.text = text
    
    # Parse alignment and set it
    p.alignment = _parse_alignment(alignment)
    
    # Set the font formatting on the run
    font = run.font
    font.size = _parse_font_size(font_size)
    font.bold = bold
    font.italic = italic
    font.name = font_name
    
    return shape

def edit_textbox(
    shape,
    text=None,
    word_wrap=None,
    font_size=None,
    bold=None,
    italic=None,
    alignment=None,
    fill_color=None,
    font_name=None
):
    """
    Edit properties of an existing textbox shape.

    :param shape: The shape object (textbox) to edit.
    :param text: New text to set. If None, leaves text unmodified.
    :param word_wrap: Boolean to enable/disable word wrap. If None, leaves unmodified.
    :param font_size: Font size (int/float or string like '12pt'). If None, leaves unmodified.
    :param bold: Boolean to set bold. If None, leaves unmodified.
    :param italic: Boolean to set italic. If None, leaves unmodified.
    :param alignment: One of 'left', 'center', 'right', 'justify'. If None, leaves unmodified.
    :param fill_color: A tuple (r, g, b) for background fill color, or None to leave unmodified.
    """

    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    # Update fill color if provided
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    # else: If you'd like to remove any existing fill if None, you could:
    # else:
    #     shape.fill.background()

    # Update word wrap if provided
    if word_wrap is not None:
        text_frame.word_wrap = word_wrap

    # If text is provided, clear existing paragraphs and add the new text
    if text is not None:
        text_frame.clear()
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = text

        # If alignment is provided, apply to the paragraph
        if alignment is not None:
            p.alignment = _parse_alignment(alignment)

        # If font formatting info is provided, apply to the run font
        font = run.font
        if font_size is not None:
            font.size = _parse_font_size(font_size)
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic

    else:
        # If no new text is given, we can selectively change existing text properties.
        for p in text_frame.paragraphs:
            if alignment is not None:
                p.alignment = _parse_alignment(alignment)
            for run in p.runs:
                font = run.font
                if font_size is not None:
                    font.size = _parse_font_size(font_size)
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if font_name is not None:
                    font.name = font_name

def add_image(slide, name, left_inch, top_inch, width_inch, height_inch, image_path):
    """
    Add an image to the slide at the specified position and size.
    
    :param slide: The slide object where the image should be placed.
    :param name: A string name/label for the shape.
    :param left_inch: Left position in inches.
    :param top_inch: Top position in inches.
    :param width_inch: Width in inches.
    :param height_inch: Height in inches.
    :param image_path: File path to the image.
    :return: The newly created picture shape object.
    """
    shape = slide.shapes.add_picture(
        image_path,
        Inches(left_inch), Inches(top_inch),
        width=Inches(width_inch), height=Inches(height_inch)
    )
    shape.name = name
    return shape

def set_shape_position(shape, left_inch, top_inch, width_inch, height_inch):
    """
    Move or resize an existing shape to the specified position/dimensions.
    
    :param shape: The shape object to be repositioned.
    :param left_inch: New left position in inches.
    :param top_inch: New top position in inches.
    :param width_inch: New width in inches.
    :param height_inch: New height in inches.
    """
    shape.left = Inches(left_inch)
    shape.top = Inches(top_inch)
    shape.width = Inches(width_inch)
    shape.height = Inches(height_inch)

def add_line_simple(slide, name, left_inch, top_inch, length_inch, thickness=2, color=(0, 0, 0), orientation="horizontal"):
    """
    Add a simple horizontal or vertical line to the slide.
    
    Parameters:
      slide: The slide object.
      name: The name/label for the line shape.
      left_inch: The left (X) coordinate in inches for the starting point.
      top_inch: The top (Y) coordinate in inches for the starting point.
      length_inch: The length of the line in inches.
      thickness: The thickness of the line in points (default is 2).
      color: An (R, G, B) tuple specifying the line color (default is black).
      orientation: "horizontal" or "vertical" (case-insensitive).
      
    Returns:
      The created line shape object.
    """
    x1 = Inches(left_inch)
    y1 = Inches(top_inch)
    
    if orientation.lower() == "horizontal":
        x2 = Inches(left_inch + length_inch)
        y2 = y1
    elif orientation.lower() == "vertical":
        x2 = x1
        y2 = Inches(top_inch + length_inch)
    else:
        raise ValueError("Orientation must be either 'horizontal' or 'vertical'")
    
    # Create a straight connector (used as a line)
    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line_shape.name = name
    
    # Set the line thickness and color
    line_shape.line.width = Pt(thickness)
    line_shape.line.color.rgb = RGBColor(*color)
    
    return line_shape

def set_paragraph_line_spacing(shape, line_spacing=1.0):
    """
    Set line spacing for all paragraphs in a textbox shape.
    E.g., line_spacing=1.5 for 1.5x spacing, 2 for double spacing, etc.
    
    :param shape: The textbox shape to modify.
    :param line_spacing: A float indicating multiple of single spacing.
    """
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.line_spacing = line_spacing  # direct float: 1.5, 2.0, etc.

def set_shape_text_margins(
    shape, 
    top_px=0, 
    right_px=0, 
    bottom_px=0, 
    left_px=0
):
    """
    Set the internal text margins (like "padding") for a textbox shape.
    python-pptx uses points or EMUs for margins, so we convert from px -> points -> EMUs as needed.
    
    Note: If your output environment uses a different PX:PT ratio, adjust _px_to_pt().
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_top = Pt(_px_to_pt(top_px))
    text_frame.margin_right = Pt(_px_to_pt(right_px))
    text_frame.margin_bottom = Pt(_px_to_pt(bottom_px))
    text_frame.margin_left = Pt(_px_to_pt(left_px))

def adjust_font_size(shape, delta=2):
    """
    Increase or decrease the current font size of all runs in a shape by `delta` points.
    If a run has no explicitly set font size (font.size is None), we can either skip it or assume a default.
    For simplicity, let's skip runs without an explicit size to avoid overwriting theme defaults.
    
    :param shape: The textbox shape to update.
    :param delta: Positive or negative integer to adjust the font size.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            current_size = run.font.size
            if current_size is not None:
                new_size = current_size.pt + delta
                # Prevent negative or zero font size
                if new_size < 1:
                    new_size = 1
                run.font.size = Pt(new_size)

def center_shape_horizontally(prs, shape):
    """
    Center a shape horizontally on the slide using the presentation's slide width.
    
    :param prs: The Presentation object (which holds slide_width).
    :param shape: The shape to center.
    """
    new_left = (prs.slide_width - shape.width) // 2
    shape.left = new_left

def center_shape_vertically(prs, shape):
    """
    Center a shape vertically on the slide using the presentation's slide height.
    
    :param prs: The Presentation object (which holds slide_height).
    :param shape: The shape to center.
    """
    new_top = (prs.slide_height - shape.height) // 2
    shape.top = new_top

def set_shape_text(shape, text, clear_first=True):
    """
    Set or replace the text of an existing shape (commonly a textbox).
    
    :param shape: The shape (textbox) whose text needs to be updated.
    :param text: The new text content.
    :param clear_first: Whether to clear existing paragraphs before adding.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    if clear_first:
        text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = text

def _set_run_font_color(run, rgb_tuple):
    """
    Manually create or replace the solidFill element in this run's XML
    to force the color if run.font.color is None or doesn't exist yet.
    """
    # Underlying run properties element
    rPr = run.font._element
    
    # Remove any existing <a:solidFill> elements to avoid duplicates
    for child in rPr.iterchildren():
        if child.tag == qn('a:solidFill'):
            rPr.remove(child)

    # Create a new solidFill element with the specified color
    solid_fill = OxmlElement('a:solidFill')
    srgb_clr = OxmlElement('a:srgbClr')
    # Format the tuple (r, g, b) into a hex string "RRGGBB"
    srgb_clr.set('val', '{:02X}{:02X}{:02X}'.format(*rgb_tuple))
    solid_fill.append(srgb_clr)
    rPr.append(solid_fill)

def set_text_style(shape, font_size=None, bold=None, italic=None, alignment=None, color=None, font_name=None):
    """
    Adjust text style on an existing textbox shape.
    
    :param shape: The textbox shape whose style is being updated.
    :param font_size: Numeric font size (e.g. 40) or None to skip.
    :param bold: Boolean or None (to skip).
    :param italic: Boolean or None (to skip).
    :param alignment: String alignment ('left', 'center', 'right', 'justify') or None (to skip).
    :param color: A tuple (r, g, b), each int from 0-255, or None (to skip).
    :param font_name: String font name (e.g., 'Arial') or None
    """
    text_frame = shape.text_frame
    # Disable auto-sizing so our manual settings are respected
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Convert the alignment string into a PP_ALIGN enum value
    parsed_alignment = _parse_alignment(alignment) if alignment else None
    
    # Convert the raw font size to a python-pptx Pt object
    parsed_font_size = _parse_font_size(font_size)

    # Iterate over paragraphs and runs in the shape
    for paragraph in text_frame.paragraphs:
        if parsed_alignment is not None:
            paragraph.alignment = parsed_alignment
        
        for run in paragraph.runs:
            # Font size
            if parsed_font_size is not None:
                run.font.size = parsed_font_size
            
            # Bold
            if bold is not None:
                run.font.bold = bold
            
            # Italic
            if italic is not None:
                run.font.italic = italic

            # Font name
            if font_name is not None:
                run.font.name = font_name
            
            # Color
            if color is not None:
                # Sometimes run.font.color may be None. We can try:
                if run.font.color is not None:
                    # If a ColorFormat object already exists, just set it
                    run.font.color.rgb = RGBColor(*color)
                else:
                    # Otherwise, manually set the run color in the underlying XML
                    _set_run_font_color(run, color)

def save_presentation(prs, file_name="poster.pptx"):
    """
    Save the current Presentation object to disk.
    
    :param prs: The Presentation object.
    :param file_name: The file path/name for the saved pptx file.
    """
    prs.save(file_name)

def set_slide_background_color(slide, rgb=(255, 255, 255)):
    """
    Sets the background color for a single Slide object.

    :param slide: A pptx.slide.Slide object
    :param rgb: A tuple of (R, G, B) color values, e.g. (255, 0, 0) for red
    """
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(*rgb)

def style_shape_border(shape, color=(30, 144, 255), thickness=2, line_style="square_dot"):
    """
    Applies a border (line) style to a given shape, where line_style is a 
    string corresponding to an MSO_LINE_DASH_STYLE enum value from python-pptx.

    Valid line_style strings (based on the doc snippet) are:
    -----------------------------------------------------------------
    'solid'        -> MSO_LINE_DASH_STYLE.SOLID
    'round_dot'    -> MSO_LINE_DASH_STYLE.ROUND_DOT
    'square_dot'   -> MSO_LINE_DASH_STYLE.SQUARE_DOT
    'dash'         -> MSO_LINE_DASH_STYLE.DASH
    'dash_dot'     -> MSO_LINE_DASH_STYLE.DASH_DOT
    'dash_dot_dot' -> MSO_LINE_DASH_STYLE.DASH_DOT_DOT
    'long_dash'    -> MSO_LINE_DASH_STYLE.LONG_DASH
    'long_dash_dot'-> MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    -----------------------------------------------------------------

    :param shape:     pptx.shapes.base.Shape object to style
    :param color:     A tuple (R, G, B) for the border color (default is (30, 144, 255))
    :param thickness: Border thickness in points (default is 2)
    :param line_style:String representing the line dash style; defaults to 'square_dot'
    """
    # Map our string keys to MSO_LINE_DASH_STYLE values from your doc snippet
    dash_style_map = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
        "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    }

    line = shape.line
    line.width = Pt(thickness)
    line.color.rgb = RGBColor(*color)

    # Default to 'solid' if the requested style isn't in dash_style_map
    dash_style_enum = dash_style_map.get(line_style.lower(), MSO_LINE_DASH_STYLE.SOLID)
    line.dash_style = dash_style_enum

def fill_textframe(shape, paragraphs_spec, vertical_anchor=None, auto_size=None):
    """
    Given an existing shape (with a text frame) and a paragraphs_spec
    describing paragraphs and runs, populate the shape's text frame.

    'paragraphs_spec' is a list of paragraphs, each containing:
      - bullet: bool
      - level: int (indent level)
      - alignment: str ("left", "center", "right", or "justify")
      - font_size: int
      - runs: list of run dictionaries, each with:
          text: str
          bold: bool
          italic: bool
          color: [r,g,b] or None
          font_size: int (optional, overrides paragraph default)
          fill_color: [r,g,b] or None

    :param vertical_anchor: Optional MSO_ANCHOR constant or string ("top", "middle", "bottom")
                           to control vertical alignment of text within textbox
    """
    text_frame = shape.text_frame
    if auto_size is None:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    elif isinstance(auto_size, str):
        auto_size_map = {
            "none": MSO_AUTO_SIZE.NONE,
            "shrink": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            "fit": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        }
        text_frame.auto_size = auto_size_map.get(auto_size.lower(), MSO_AUTO_SIZE.NONE)
    else:
        text_frame.auto_size = auto_size
    text_frame.word_wrap = True

    # Set vertical anchor if provided
    if vertical_anchor is not None:
        if isinstance(vertical_anchor, str):
            anchor_map = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            text_frame.vertical_anchor = anchor_map.get(vertical_anchor.lower(), MSO_ANCHOR.TOP)
        else:
            text_frame.vertical_anchor = vertical_anchor

    # Clear out existing paragraphs
    text_frame.clear()

    for p_data in paragraphs_spec:
        p = text_frame.add_paragraph()
        
        # # bulleting
        # p.bullet = p_data.get("bullet", False)
        
        # bullet level (indent)
        p.level = p_data.get("level", 0)
        
        # paragraph alignment
        align_str = p_data.get("alignment", "left")
        p.alignment = _parse_alignment(align_str)
        
        # paragraph-level font size
        default_font_size = p_data.get("font_size", 24)
        p.font.size = Pt(default_font_size)

        # Add runs
        runs_spec = p_data.get("runs", [])
        for run_info in runs_spec:
            run = p.add_run()
            if p_data.get("bullet", False):
                if p.level == 0:
                    run.text = '\u2022' + run_info.get("text", "")
                elif p.level == 1:
                    run.text = '\u25E6' + run_info.get("text", "")
                else:
                    run.text = '\u25AA' + run_info.get("text", "")
            else:
                run.text = run_info.get("text", "")

            # Font styling
            font = run.font
            font.bold = run_info.get("bold", False)
            font.italic = run_info.get("italic", False)

            # If run-specific color was provided
            color_tuple = run_info.get("color", None)
            if (
                color_tuple
                and len(color_tuple) == 3
                and all(isinstance(c, int) for c in color_tuple)
            ):
                if run.font.color is not None:
                    # If a ColorFormat object already exists, just set it
                    run.font.color.rgb = RGBColor(*color_tuple)
                else:
                    # Otherwise, manually set the run color in the underlying XML
                    _set_run_font_color(run, color_tuple)

            # If run-specific font size was provided
            if "font_size" in run_info:
                font.size = Pt(run_info["font_size"])

            # If run-specific shape fill color was provided:
            fill_color_tuple = run_info.get("fill_color", None)
            if (
                fill_color_tuple
                and len(fill_color_tuple) == 3
                and all(isinstance(c, int) for c in fill_color_tuple)
            ):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*fill_color_tuple)


def add_border_hierarchy(
    prs,
    name_to_hierarchy: dict,
    hierarchy: int,
    border_color=RGBColor(255, 0, 0),
    border_width=2,
    fill_boxes: bool = False,
    fill_color=RGBColor(255, 0, 0),
    regardless=False
):
    """
    Iterates over all slides and shapes in the Presentation 'prs'.
    - For shapes whose name maps to the given 'hierarchy' in 'name_to_hierarchy' (or if 'regardless'
      is True), draws a red border. Optionally fills the shape with red if 'fill_boxes' is True.
    - For all other shapes, removes their border and hides any text.

    Returns:
        labeled_elements: dict of shape geometry for ALL shapes, regardless of hierarchy match.
    """
    border_width = Pt(border_width)
    labeled_elements = {}

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Record basic geometry in labeled_elements
            shape_name = shape.name if hasattr(shape, 'name') else f"Shape_{slide_idx}_{shape_idx}"
            labeled_elements[shape_name] = {
                'left': f"{emu_to_inches(shape.left):.2f} Inches",
                'top': f"{emu_to_inches(shape.top):.2f} Inches",
                'width': f"{emu_to_inches(shape.width):.2f} Inches",
                'height': f"{emu_to_inches(shape.height):.2f} Inches",
            }

            # Determine if this shape should have a border
            current_hierarchy = name_to_hierarchy.get(shape_name, None)
            if current_hierarchy is None:
                # Optional: Print a debug message if the shape’s name isn’t in the dict
                print(f"Warning: shape '{shape_name}' not found in name_to_hierarchy.")

            try:
                if current_hierarchy == hierarchy or regardless:
                    # Draw border
                    shape.line.fill.solid()
                    shape.line.fill.fore_color.rgb = border_color
                    shape.line.width = border_width

                    # Optionally fill the shape with red color
                    if fill_boxes:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = fill_color
                else:
                    # Remove border
                    shape.line.width = Pt(0)
                    shape.line.fill.background()

                    # Hide text if present
                    if shape.has_text_frame:
                        shape.text_frame.text = ""
            except Exception as e:
                print(f"Could not process shape '{shape_name}' (type={shape.shape_type}): {e}")

    return labeled_elements


def get_visual_cues(name_to_hierarchy, identifier, poster_path='poster.pptx'):
    prs = pptx.Presentation(poster_path)

    position_dict_1 = add_border_hierarchy(prs, name_to_hierarchy, 1, border_width=10)
    json.dump(position_dict_1, open(f"tmp/position_dict_1_<{identifier}>.json", "w"))

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_1.pptx")

    prs = pptx.Presentation(poster_path)

    add_border_hierarchy(prs, name_to_hierarchy, 1, border_width=10, fill_boxes=True)
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_1_filled.pptx")

    prs = pptx.Presentation(poster_path)

    position_dict_2 = add_border_hierarchy(prs, name_to_hierarchy, 2, border_width=10)
    json.dump(position_dict_2, open(f"tmp/position_dict_2_<{identifier}>.json", "w"))

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_2.pptx")

    prs = pptx.Presentation(poster_path)

    add_border_hierarchy(prs, name_to_hierarchy, 2, border_width=10, fill_boxes=True)

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_2_filled.pptx")

'''


documentation = r'''
create_poster(width_inch=48, height_inch=36):
    """
    Create a new Presentation object, set its slide size (e.g., 48x36 inches).
    
    :param width_inch: Float or int specifying width in inches (default 48).
    :param height_inch: Float or int specifying height in inches (default 36).
    :return: A python-pptx Presentation object.
    """
    
add_blank_slide(prs):
    """
    Add a blank slide to the Presentation (layout index 6 is typically blank).
    
    :param prs: The Presentation object to add a slide to.
    :return: The newly added slide object.
    """

def shape_fill_color(shape, fill_color):
    """
    Set the fill color of a shape to the specified RGB color.

    :param shape: The shape object to modify.
    :param fill_color: A tuple (r, g, b) for the fill color.
    """

def add_textbox(
    slide, 
    name, 
    left_inch, 
    top_inch, 
    width_inch, 
    height_inch, 
    text="", 
    word_wrap=True,
    font_size=40,
    bold=False,
    italic=False,
    alignment="left",
    fill_color=None,
    font_name="Arial"
):
    """
    Create a textbox shape on the given slide, optionally fill its background with
    a color if fill_color is specified as (r, g, b).
    
    :param slide: Slide object to place the textbox on.
    :param name: Name for the shape (shape.name).
    :param left_inch: Left coordinate (in inches).
    :param top_inch: Top coordinate (in inches).
    :param width_inch: Width (in inches).
    :param height_inch: Height (in inches).
    :param text: Text to display in the textbox.
    :param word_wrap: If True, wrap text in the textbox.
    :param font_size: Numeric font size (e.g. 40).
    :param bold: Boolean to set run.font.bold.
    :param italic: Boolean to set run.font.italic.
    :param alignment: String alignment: "left", "center", "right", or "justify".
    :param fill_color: (r, g, b) tuple for solid fill background color, or None to skip.
    :param font_name: String font name (e.g., "Arial").
    :return: The newly created textbox shape.
    """

add_image(slide, name, left_inch, top_inch, width_inch, height_inch, image_path):
    """
    Add an image to the slide at the specified position and size.
    
    :param slide: The slide object where the image should be placed.
    :param name: A string name/label for the shape.
    :param left_inch: Left position in inches.
    :param top_inch: Top position in inches.
    :param width_inch: Width in inches.
    :param height_inch: Height in inches.
    :param image_path: File path to the image.
    :return: The newly created picture shape object.
    """

set_shape_position(shape, left_inch, top_inch, width_inch, height_inch):
    """
    Move or resize an existing shape to the specified position/dimensions.
    
    :param shape: The shape object to be repositioned.
    :param left_inch: New left position in inches.
    :param top_inch: New top position in inches.
    :param width_inch: New width in inches.
    :param height_inch: New height in inches.
    """

def set_text_style(shape, font_size=None, bold=None, italic=None, alignment=None, color=None, font_name=None):
    """
    Adjust text style on an existing textbox shape.
    
    :param shape: The textbox shape whose style is being updated.
    :param font_size: Numeric font size (e.g. 40) or None to skip.
    :param bold: Boolean or None (to skip).
    :param italic: Boolean or None (to skip).
    :param alignment: String alignment ('left', 'center', 'right', 'justify') or None (to skip).
    :param color: A tuple (r, g, b), each int from 0-255, or None (to skip).
    :param font_name: String font name (e.g., 'Arial') or None
    """

add_line_simple(slide, name, left_inch, top_inch, length_inch, thickness=2, color=(0, 0, 0), orientation="horizontal"):
    """
    Add a simple horizontal or vertical line to the slide.
    
    Parameters:
      slide: The slide object.
      name: The name/label for the line shape.
      left_inch: The left (X) coordinate in inches for the starting point.
      top_inch: The top (Y) coordinate in inches for the starting point.
      length_inch: The length of the line in inches.
      thickness: The thickness of the line in points (default is 2).
      color: An (R, G, B) tuple specifying the line color (default is black).
      orientation: "horizontal" or "vertical" (case-insensitive).
      
    Returns:
      The created line shape object.
    """

set_paragraph_line_spacing(shape, line_spacing=1.0):
    """
    Set line spacing for all paragraphs in a textbox shape.
    E.g., line_spacing=1.5 for 1.5x spacing, 2 for double spacing, etc.
    
    :param shape: The textbox shape to modify.
    :param line_spacing: A float indicating multiple of single spacing.
    """

set_shape_text_margins(
    shape, 
    top_px=0, 
    right_px=0, 
    bottom_px=0, 
    left_px=0
):
    """
    Set the internal text margins (like "padding") for a textbox shape.
    python-pptx uses points or EMUs for margins, so we convert from px -> points -> EMUs as needed.
    
    Note: If your output environment uses a different PX:PT ratio, adjust _px_to_pt().
    """

adjust_font_size(shape, delta=2):
    """
    Increase or decrease the current font size of all runs in a shape by `delta` points.
    If a run has no explicitly set font size (font.size is None), we can either skip it or assume a default.
    For simplicity, let's skip runs without an explicit size to avoid overwriting theme defaults.
    
    :param shape: The textbox shape to update.
    :param delta: Positive or negative integer to adjust the font size.
    """

def set_slide_background_color(slide, rgb=(255, 255, 255)):
    """
    Sets the background color for a single Slide object.

    :param slide: A pptx.slide.Slide object
    :param rgb: A tuple of (R, G, B) color values, e.g. (255, 0, 0) for red
    """

def style_shape_border(shape, color=(30, 144, 255), thickness=2, line_style="square_dot"):
    """
    Applies a border (line) style to a given shape, where line_style is a 
    string corresponding to an MSO_LINE_DASH_STYLE enum value from python-pptx.

    Valid line_style strings (based on the doc snippet) are:
    -----------------------------------------------------------------
    'solid'        -> MSO_LINE_DASH_STYLE.SOLID
    'round_dot'    -> MSO_LINE_DASH_STYLE.ROUND_DOT
    'square_dot'   -> MSO_LINE_DASH_STYLE.SQUARE_DOT
    'dash'         -> MSO_LINE_DASH_STYLE.DASH
    'dash_dot'     -> MSO_LINE_DASH_STYLE.DASH_DOT
    'dash_dot_dot' -> MSO_LINE_DASH_STYLE.DASH_DOT_DOT
    'long_dash'    -> MSO_LINE_DASH_STYLE.LONG_DASH
    'long_dash_dot'-> MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    -----------------------------------------------------------------

    :param shape:     pptx.shapes.base.Shape object to style
    :param color:     A tuple (R, G, B) for the border color (default is (30, 144, 255))
    :param thickness: Border thickness in points (default is 2)
    :param line_style:String representing the line dash style; defaults to 'square_dot'
    """

save_presentation(prs, file_name="poster.pptx"):
    """
    Save the current Presentation object to disk.
    
    :param prs: The Presentation object.
    :param file_name: The file path/name for the saved pptx file.
    """

--------------------------------------

Example usage:
poster = create_poster(width_inch=48, height_inch=36)
slide = add_blank_slide(poster)
# Set this particular slide's background to light gray
set_slide_background_color(slide, (200, 200, 200))

title_text_box = add_textbox(
    slide, 
    name='title', 
    left_inch=5, 
    top_inch=0, 
    width_inch=30, 
    height_inch=5, 
    text="Poster Title", 
    word_wrap=True,
    font_size=100,
    bold=True,
    italic=False,
    alignment="center",
    fill_color=(255, 255, 255),  # Fill color
    font_name="Arial"
)

shape_fill_color(title_text_box, fill_color=(173, 216, 230)) # Fill color

# Apply a dashed border with "square_dot"
style_shape_border(title_text_box, color=(30, 144, 255), thickness=8, line_style="square_dot")
image = add_image(slide, 'img', 10, 25, 30, 30, 'data/poster_exp/pdf/attention/_page_3_Figure_0.jpeg')

set_shape_position(image, 10, 25, 15, 15)
set_shape_position(image, 10, 5, 20, 15)

set_text_style(title_text_box, font_size=60, bold=True, italic=True, alignment='center', color=(255, 0, 0), font_name='Times New Roman')

added_line = add_line_simple(
    slide,
    'separation_line',
    20,
    0,
    20,
    thickness=2,   # in points
    color=(120, 120, 20),
    orientation='vertical'
)

set_shape_text_margins(
    title_text_box, 
    top_px=10, 
    right_px=20, 
    bottom_px=30, 
    left_px=40
)

adjust_font_size(title_text_box, delta=-20)

set_paragraph_line_spacing(title_text_box, line_spacing=2.0)

save_presentation(poster, file_name="poster.pptx")

'''


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import pptx

from pptx.enum.text import MSO_AUTO_SIZE

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def _px_to_pt(px):
    """
    Approximate conversion from pixels to points.
    A common assumption is 1px ~ 0.75pt.
    Adjust as needed for your environment.
    """
    return px * 0.75

def _parse_font_size(font_size):
    """
    Internal helper to convert a numeric font size (e.g., 12) 
    to a python-pptx Pt object. If it's already a Pt, return as-is.
    """
    if font_size is None:
        return None
    if isinstance(font_size, (int, float)):
        return Pt(font_size)
    return font_size  # Assume user provided a Pt object already

def _parse_alignment(alignment):
    """
    Internal helper to convert a string alignment (e.g., "left", "center") 
    to the corresponding PP_ALIGN constant. 
    Default to PP_ALIGN.LEFT if unrecognized or None.
    """
    if not isinstance(alignment, str):
        # If user passed None or something else, default to PP_ALIGN.LEFT
        return PP_ALIGN.LEFT

    alignment = alignment.lower().strip()
    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return alignment_map.get(alignment, PP_ALIGN.LEFT)

def create_poster(width_inch=48, height_inch=36):
    """
    Create a new Presentation object, set its slide size (e.g., 48x36 inches).
    
    :param width_inch: Float or int specifying width in inches (default 48).
    :param height_inch: Float or int specifying height in inches (default 36).
    :return: A python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = Inches(width_inch)
    prs.slide_height = Inches(height_inch)
    return prs

def add_blank_slide(prs):
    """
    Add a blank slide to the Presentation (layout index 6 is typically blank).
    
    :param prs: The Presentation object to add a slide to.
    :return: The newly added slide object.
    """
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def shape_fill_color(shape, fill_color):
    """
    Set the fill color of a shape to the specified RGB color.

    :param shape: The shape object to modify.
    :param fill_color: A tuple (r, g, b) for the fill color.
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)

def add_textbox(
    slide, 
    name, 
    left_inch, 
    top_inch, 
    width_inch, 
    height_inch, 
    text="", 
    word_wrap=True,
    font_size=40,
    bold=False,
    italic=False,
    alignment="left",
    fill_color=None,
    font_name="Arial"
):
    """
    Create a textbox shape on the given slide, optionally fill its background with
    a color if fill_color is specified as (r, g, b).
    
    :param slide: Slide object to place the textbox on.
    :param name: Name for the shape (shape.name).
    :param left_inch: Left coordinate (in inches).
    :param top_inch: Top coordinate (in inches).
    :param width_inch: Width (in inches).
    :param height_inch: Height (in inches).
    :param text: Text to display in the textbox.
    :param word_wrap: If True, wrap text in the textbox.
    :param font_size: Numeric font size (e.g. 40).
    :param bold: Boolean to set run.font.bold.
    :param italic: Boolean to set run.font.italic.
    :param alignment: String alignment: "left", "center", "right", or "justify".
    :param fill_color: (r, g, b) tuple for solid fill background color, or None to skip.
    :param font_name: String font name (e.g., "Arial").
    :return: The newly created textbox shape.
    """
    shape = slide.shapes.add_textbox(
        Inches(left_inch), Inches(top_inch),
        Inches(width_inch), Inches(height_inch)
    )
    
    shape.name = name
    
    # If a fill color is specified, apply a solid fill
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        # Otherwise, set "no fill" if you want it transparent
        shape.fill.background()

    text_frame = shape.text_frame
    # Turn off auto-size to ensure stable font size, etc.
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.word_wrap = word_wrap
    
    # Clear any default paragraphs
    text_frame.clear()
    
    # Add a new paragraph
    p = text_frame.add_paragraph()
    # Instead of setting p.text, explicitly create a Run
    run = p.add_run()
    run.text = text
    
    # Parse alignment and set it
    p.alignment = _parse_alignment(alignment)
    
    # Set the font formatting on the run
    font = run.font
    font.size = _parse_font_size(font_size)
    font.bold = bold
    font.italic = italic
    font.name = font_name
    
    return shape

def fill_textframe(shape, paragraphs_spec, vertical_anchor=None, auto_size=None):
    """
    Given an existing shape (with a text frame) and a paragraphs_spec
    describing paragraphs and runs, populate the shape's text frame.

    'paragraphs_spec' is a list of paragraphs, each containing:
      - bullet: bool
      - level: int (indent level)
      - alignment: str ("left", "center", "right", or "justify")
      - font_size: int
      - runs: list of run dictionaries, each with:
          text: str
          bold: bool
          italic: bool
          color: [r,g,b] or None
          font_size: int (optional, overrides paragraph default)
          fill_color: [r,g,b] or None

    :param vertical_anchor: Optional MSO_ANCHOR constant or string ("top", "middle", "bottom")
                           to control vertical alignment of text within textbox
    """
    text_frame = shape.text_frame
    if auto_size is None:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    elif isinstance(auto_size, str):
        auto_size_map = {
            "none": MSO_AUTO_SIZE.NONE,
            "shrink": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            "fit": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        }
        text_frame.auto_size = auto_size_map.get(auto_size.lower(), MSO_AUTO_SIZE.NONE)
    else:
        text_frame.auto_size = auto_size
    text_frame.word_wrap = True

    # Set vertical anchor if provided
    if vertical_anchor is not None:
        if isinstance(vertical_anchor, str):
            anchor_map = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            text_frame.vertical_anchor = anchor_map.get(vertical_anchor.lower(), MSO_ANCHOR.TOP)
        else:
            text_frame.vertical_anchor = vertical_anchor

    # Clear out existing paragraphs
    text_frame.clear()

    for p_data in paragraphs_spec:
        p = text_frame.add_paragraph()
        
        # # bulleting
        # p.bullet = p_data.get("bullet", False)
        
        # bullet level (indent)
        p.level = p_data.get("level", 0)
        
        # paragraph alignment
        align_str = p_data.get("alignment", "left")
        p.alignment = _parse_alignment(align_str)
        
        # paragraph-level font size
        default_font_size = p_data.get("font_size", 24)
        p.font.size = Pt(default_font_size)

        # Add runs
        runs_spec = p_data.get("runs", [])
        for run_info in runs_spec:
            run = p.add_run()
            if p_data.get("bullet", False):
                if p.level == 0:
                    run.text = '\u2022' + run_info.get("text", "")
                elif p.level == 1:
                    run.text = '\u25E6' + run_info.get("text", "")
                else:
                    run.text = '\u25AA' + run_info.get("text", "")
            else:
                run.text = run_info.get("text", "")

            # Font styling
            font = run.font
            font.bold = run_info.get("bold", False)
            font.italic = run_info.get("italic", False)

            # If run-specific color was provided
            color_tuple = run_info.get("color", None)
            if (
                color_tuple
                and len(color_tuple) == 3
                and all(isinstance(c, int) for c in color_tuple)
            ):
                if run.font.color is not None:
                    # If a ColorFormat object already exists, just set it
                    run.font.color.rgb = RGBColor(*color_tuple)
                else:
                    # Otherwise, manually set the run color in the underlying XML
                    _set_run_font_color(run, color_tuple)

            # If run-specific font size was provided
            if "font_size" in run_info:
                font.size = Pt(run_info["font_size"])

            # If run-specific shape fill color was provided:
            fill_color_tuple = run_info.get("fill_color", None)
            if (
                fill_color_tuple
                and len(fill_color_tuple) == 3
                and all(isinstance(c, int) for c in fill_color_tuple)
            ):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*fill_color_tuple)


def edit_textbox(
    shape,
    text=None,
    word_wrap=None,
    font_size=None,
    bold=None,
    italic=None,
    alignment=None,
    fill_color=None,
    font_name=None
):
    """
    Edit properties of an existing textbox shape.

    :param shape: The shape object (textbox) to edit.
    :param text: New text to set. If None, leaves text unmodified.
    :param word_wrap: Boolean to enable/disable word wrap. If None, leaves unmodified.
    :param font_size: Font size (int/float or string like '12pt'). If None, leaves unmodified.
    :param bold: Boolean to set bold. If None, leaves unmodified.
    :param italic: Boolean to set italic. If None, leaves unmodified.
    :param alignment: One of 'left', 'center', 'right', 'justify'. If None, leaves unmodified.
    :param fill_color: A tuple (r, g, b) for background fill color, or None to leave unmodified.
    """

    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    # Update fill color if provided
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    # else: If you'd like to remove any existing fill if None, you could:
    # else:
    #     shape.fill.background()

    # Update word wrap if provided
    if word_wrap is not None:
        text_frame.word_wrap = word_wrap

    # If text is provided, clear existing paragraphs and add the new text
    if text is not None:
        text_frame.clear()
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = text

        # If alignment is provided, apply to the paragraph
        if alignment is not None:
            p.alignment = _parse_alignment(alignment)

        # If font formatting info is provided, apply to the run font
        font = run.font
        if font_size is not None:
            font.size = _parse_font_size(font_size)
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic

    else:
        # If no new text is given, we can selectively change existing text properties.
        for p in text_frame.paragraphs:
            if alignment is not None:
                p.alignment = _parse_alignment(alignment)
            for run in p.runs:
                font = run.font
                if font_size is not None:
                    font.size = _parse_font_size(font_size)
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if font_name is not None:
                    font.name = font_name

def add_image(slide, name, left_inch, top_inch, width_inch, height_inch, image_path):
    """
    Add an image to the slide at the specified position and size.
    
    :param slide: The slide object where the image should be placed.
    :param name: A string name/label for the shape.
    :param left_inch: Left position in inches.
    :param top_inch: Top position in inches.
    :param width_inch: Width in inches.
    :param height_inch: Height in inches.
    :param image_path: File path to the image.
    :return: The newly created picture shape object.
    """
    shape = slide.shapes.add_picture(
        image_path,
        Inches(left_inch), Inches(top_inch),
        width=Inches(width_inch), height=Inches(height_inch)
    )
    shape.name = name
    return shape

def set_shape_position(shape, left_inch, top_inch, width_inch, height_inch):
    """
    Move or resize an existing shape to the specified position/dimensions.
    
    :param shape: The shape object to be repositioned.
    :param left_inch: New left position in inches.
    :param top_inch: New top position in inches.
    :param width_inch: New width in inches.
    :param height_inch: New height in inches.
    """
    shape.left = Inches(left_inch)
    shape.top = Inches(top_inch)
    shape.width = Inches(width_inch)
    shape.height = Inches(height_inch)

def add_line_simple(slide, name, left_inch, top_inch, length_inch, thickness=2, color=(0, 0, 0), orientation="horizontal"):
    """
    Add a simple horizontal or vertical line to the slide.
    
    Parameters:
      slide: The slide object.
      name: The name/label for the line shape.
      left_inch: The left (X) coordinate in inches for the starting point.
      top_inch: The top (Y) coordinate in inches for the starting point.
      length_inch: The length of the line in inches.
      thickness: The thickness of the line in points (default is 2).
      color: An (R, G, B) tuple specifying the line color (default is black).
      orientation: "horizontal" or "vertical" (case-insensitive).
      
    Returns:
      The created line shape object.
    """
    x1 = Inches(left_inch)
    y1 = Inches(top_inch)
    
    if orientation.lower() == "horizontal":
        x2 = Inches(left_inch + length_inch)
        y2 = y1
    elif orientation.lower() == "vertical":
        x2 = x1
        y2 = Inches(top_inch + length_inch)
    else:
        raise ValueError("Orientation must be either 'horizontal' or 'vertical'")
    
    # Create a straight connector (used as a line)
    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line_shape.name = name
    
    # Set the line thickness and color
    line_shape.line.width = Pt(thickness)
    line_shape.line.color.rgb = RGBColor(*color)
    
    return line_shape

def set_paragraph_line_spacing(shape, line_spacing=1.0):
    """
    Set line spacing for all paragraphs in a textbox shape.
    E.g., line_spacing=1.5 for 1.5x spacing, 2 for double spacing, etc.
    
    :param shape: The textbox shape to modify.
    :param line_spacing: A float indicating multiple of single spacing.
    """
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.line_spacing = line_spacing  # direct float: 1.5, 2.0, etc.

def set_shape_text_margins(
    shape, 
    top_px=0, 
    right_px=0, 
    bottom_px=0, 
    left_px=0
):
    """
    Set the internal text margins (like "padding") for a textbox shape.
    python-pptx uses points or EMUs for margins, so we convert from px -> points -> EMUs as needed.
    
    Note: If your output environment uses a different PX:PT ratio, adjust _px_to_pt().
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_top = Pt(_px_to_pt(top_px))
    text_frame.margin_right = Pt(_px_to_pt(right_px))
    text_frame.margin_bottom = Pt(_px_to_pt(bottom_px))
    text_frame.margin_left = Pt(_px_to_pt(left_px))

def adjust_font_size(shape, delta=2):
    """
    Increase or decrease the current font size of all runs in a shape by `delta` points.
    If a run has no explicitly set font size (font.size is None), we can either skip it or assume a default.
    For simplicity, let's skip runs without an explicit size to avoid overwriting theme defaults.
    
    :param shape: The textbox shape to update.
    :param delta: Positive or negative integer to adjust the font size.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            current_size = run.font.size
            if current_size is not None:
                new_size = current_size.pt + delta
                # Prevent negative or zero font size
                if new_size < 1:
                    new_size = 1
                run.font.size = Pt(new_size)

def center_shape_horizontally(prs, shape):
    """
    Center a shape horizontally on the slide using the presentation's slide width.
    
    :param prs: The Presentation object (which holds slide_width).
    :param shape: The shape to center.
    """
    new_left = (prs.slide_width - shape.width) // 2
    shape.left = new_left

def center_shape_vertically(prs, shape):
    """
    Center a shape vertically on the slide using the presentation's slide height.
    
    :param prs: The Presentation object (which holds slide_height).
    :param shape: The shape to center.
    """
    new_top = (prs.slide_height - shape.height) // 2
    shape.top = new_top

def set_shape_text(shape, text, clear_first=True):
    """
    Set or replace the text of an existing shape (commonly a textbox).
    
    :param shape: The shape (textbox) whose text needs to be updated.
    :param text: The new text content.
    :param clear_first: Whether to clear existing paragraphs before adding.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    if clear_first:
        text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = text

def _set_run_font_color(run, rgb_tuple):
    """
    Manually create or replace the solidFill element in this run's XML
    to force the color if run.font.color is None or doesn't exist yet.
    """
    # Underlying run properties element
    rPr = run.font._element
    
    # Remove any existing <a:solidFill> elements to avoid duplicates
    for child in rPr.iterchildren():
        if child.tag == qn('a:solidFill'):
            rPr.remove(child)

    # Create a new solidFill element with the specified color
    solid_fill = OxmlElement('a:solidFill')
    srgb_clr = OxmlElement('a:srgbClr')
    # Format the tuple (r, g, b) into a hex string "RRGGBB"
    srgb_clr.set('val', '{:02X}{:02X}{:02X}'.format(*rgb_tuple))
    solid_fill.append(srgb_clr)
    rPr.append(solid_fill)

def set_text_style(shape, font_size=None, bold=None, italic=None, alignment=None, color=None, font_name=None):
    """
    Adjust text style on an existing textbox shape.
    
    :param shape: The textbox shape whose style is being updated.
    :param font_size: Numeric font size (e.g. 40) or None to skip.
    :param bold: Boolean or None (to skip).
    :param italic: Boolean or None (to skip).
    :param alignment: String alignment ('left', 'center', 'right', 'justify') or None (to skip).
    :param color: A tuple (r, g, b), each int from 0-255, or None (to skip).
    :param font_name: String font name (e.g., 'Arial') or None
    """
    text_frame = shape.text_frame
    # Disable auto-sizing so our manual settings are respected
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Convert the alignment string into a PP_ALIGN enum value
    parsed_alignment = _parse_alignment(alignment) if alignment else None
    
    # Convert the raw font size to a python-pptx Pt object
    parsed_font_size = _parse_font_size(font_size)

    # Iterate over paragraphs and runs in the shape
    for paragraph in text_frame.paragraphs:
        if parsed_alignment is not None:
            paragraph.alignment = parsed_alignment
        
        for run in paragraph.runs:
            # Font size
            if parsed_font_size is not None:
                run.font.size = parsed_font_size
            
            # Bold
            if bold is not None:
                run.font.bold = bold
            
            # Italic
            if italic is not None:
                run.font.italic = italic

            # Font name
            if font_name is not None:
                run.font.name = font_name
            
            # Color
            if color is not None:
                # Sometimes run.font.color may be None. We can try:
                if run.font.color is not None:
                    # If a ColorFormat object already exists, just set it
                    run.font.color.rgb = RGBColor(*color)
                else:
                    # Otherwise, manually set the run color in the underlying XML
                    _set_run_font_color(run, color)

def save_presentation(prs, file_name="poster.pptx"):
    """
    Save the current Presentation object to disk.
    
    :param prs: The Presentation object.
    :param file_name: The file path/name for the saved pptx file.
    """
    prs.save(file_name)

def set_slide_background_color(slide, rgb=(255, 255, 255)):
    """
    Sets the background color for a single Slide object.

    :param slide: A pptx.slide.Slide object
    :param rgb: A tuple of (R, G, B) color values, e.g. (255, 0, 0) for red
    """
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(*rgb)

def style_shape_border(shape, color=(30, 144, 255), thickness=2, line_style="square_dot"):
    """
    Applies a border (line) style to a given shape, where line_style is a 
    string corresponding to an MSO_LINE_DASH_STYLE enum value from python-pptx.

    Valid line_style strings (based on the doc snippet) are:
    -----------------------------------------------------------------
    'solid'        -> MSO_LINE_DASH_STYLE.SOLID
    'round_dot'    -> MSO_LINE_DASH_STYLE.ROUND_DOT
    'square_dot'   -> MSO_LINE_DASH_STYLE.SQUARE_DOT
    'dash'         -> MSO_LINE_DASH_STYLE.DASH
    'dash_dot'     -> MSO_LINE_DASH_STYLE.DASH_DOT
    'dash_dot_dot' -> MSO_LINE_DASH_STYLE.DASH_DOT_DOT
    'long_dash'    -> MSO_LINE_DASH_STYLE.LONG_DASH
    'long_dash_dot'-> MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    -----------------------------------------------------------------

    :param shape:     pptx.shapes.base.Shape object to style
    :param color:     A tuple (R, G, B) for the border color (default is (30, 144, 255))
    :param thickness: Border thickness in points (default is 2)
    :param line_style:String representing the line dash style; defaults to 'square_dot'
    """
    # Map our string keys to MSO_LINE_DASH_STYLE values from your doc snippet
    dash_style_map = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
        "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT
    }

    line = shape.line
    line.width = Pt(thickness)
    line.color.rgb = RGBColor(*color)

    # Default to 'solid' if the requested style isn't in dash_style_map
    dash_style_enum = dash_style_map.get(line_style.lower(), MSO_LINE_DASH_STYLE.SOLID)
    line.dash_style = dash_style_enum

def get_visual_cues(name_to_hierarchy, identifier, poster_path='poster.pptx'):
    prs = pptx.Presentation(poster_path)

    position_dict_1 = add_border_hierarchy(prs, name_to_hierarchy, 1, border_width=10)
    json.dump(position_dict_1, open(f"tmp/position_dict_1_<{identifier}>.json", "w"))

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_1.pptx")

    prs = pptx.Presentation(poster_path)

    add_border_hierarchy(prs, name_to_hierarchy, 1, border_width=10, fill_boxes=True)
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_1_filled.pptx")

    prs = pptx.Presentation(poster_path)

    position_dict_2 = add_border_hierarchy(prs, name_to_hierarchy, 2, border_width=10)
    json.dump(position_dict_2, open(f"tmp/position_dict_2_<{identifier}>.json", "w"))

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_2.pptx")

    prs = pptx.Presentation(poster_path)

    add_border_hierarchy(prs, name_to_hierarchy, 2, border_width=10, fill_boxes=True)

    # Save the presentation to disk.
    save_presentation(prs, file_name=f"tmp/poster_<{identifier}>_hierarchy_2_filled.pptx")

from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def add_border(
    prs,
    border_color=RGBColor(255, 0, 0),   # Red border for shapes
    border_width=Pt(2),                # 2-point border width
):
    """
    Iterates over all slides and shapes in the Presentation 'prs', applies a 
    red border to each shape, and places a transparent (no fill).

    Args:
        prs: The Presentation object to modify.
        border_color: RGBColor for the shape border color (default: red).
        border_width: The width of the shape border (Pt).
    """
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                # --- 1) Add red border to the shape (if supported) ---
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                if hasattr(shape, 'name'):
                    labeled_elements[shape.name] = {
                        'left': f'{emu_to_inches(shape.left)} Inches',
                        'top': f'{emu_to_inches(shape.top)} Inches',
                        'width': f'{emu_to_inches(shape.width)} Inches',
                        'height': f'{emu_to_inches(shape.height)} Inches',
                    }

            except Exception as e:
                # If the shape doesn't support borders or text, skip gracefully
                print(f"Could not add border to shape (type={shape.shape_type}): {e}")
    
    return labeled_elements

def get_hierarchy(outline, hierarchy=1):
    name_to_hierarchy = {}
    for key, section in outline.items():
        if key == "meta":
            continue
        name_to_hierarchy[section['name']] = hierarchy
        if 'subsections' in section:
            name_to_hierarchy.update(get_hierarchy(section['subsections'], hierarchy+1))
    return name_to_hierarchy

def get_hierarchy_by_keys(outline, hierarchy=1):
    name_to_hierarchy = {}
    for key, section in outline.items():
        if key == "meta":
            continue
        name_to_hierarchy[key] = hierarchy
        if 'subsections' in section:
            name_to_hierarchy.update(get_hierarchy_by_keys(section['subsections'], hierarchy+1))
    return name_to_hierarchy

def rename_keys_with_name(data):
    """
    Recursively rename dictionary keys to data['name'] if:
      - The value is a dict,
      - It contains a 'name' field.
    Otherwise, keep the original key.
    """
    if not isinstance(data, dict):
        # If it's not a dictionary (e.g. list or scalar), just return it as-is
        return data

    new_dict = {}
    for key, value in data.items():
        if isinstance(value, dict) and "name" in value:
            # Rename the key to whatever 'name' is in the nested dictionary
            new_key = value["name"]
            # Recursively process the value (which may contain its own subsections)
            new_dict[new_key] = rename_keys_with_name(value)
        else:
            # Keep the same key if there's no 'name' in value or it's not a dictionary
            new_dict[key] = rename_keys_with_name(value)

    return new_dict

def add_border_hierarchy(
    prs,
    name_to_hierarchy: dict,
    hierarchy: int,
    border_color=RGBColor(255, 0, 0),
    border_width=2,
    fill_boxes: bool = False,
    fill_color=RGBColor(255, 0, 0),
    regardless=False
):
    """
    Iterates over all slides and shapes in the Presentation 'prs'.
    - For shapes whose name maps to the given 'hierarchy' in 'name_to_hierarchy' (or if 'regardless'
      is True), draws a red border. Optionally fills the shape with red if 'fill_boxes' is True.
    - For all other shapes, removes their border and hides any text.

    Returns:
        labeled_elements: dict of shape geometry for ALL shapes, regardless of hierarchy match.
    """
    border_width = Pt(border_width)
    labeled_elements = {}

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Record basic geometry in labeled_elements
            shape_name = shape.name if hasattr(shape, 'name') else f"Shape_{slide_idx}_{shape_idx}"
            labeled_elements[shape_name] = {
                'left': f"{emu_to_inches(shape.left):.2f} Inches",
                'top': f"{emu_to_inches(shape.top):.2f} Inches",
                'width': f"{emu_to_inches(shape.width):.2f} Inches",
                'height': f"{emu_to_inches(shape.height):.2f} Inches",
            }

            # Determine if this shape should have a border
            current_hierarchy = name_to_hierarchy.get(shape_name, None)
            if current_hierarchy is None:
                # Optional: Print a debug message if the shape’s name isn’t in the dict
                print(f"Warning: shape '{shape_name}' not found in name_to_hierarchy.")

            try:
                if current_hierarchy == hierarchy or regardless:
                    # Draw border
                    shape.line.fill.solid()
                    shape.line.fill.fore_color.rgb = border_color
                    shape.line.width = border_width

                    # Optionally fill the shape with red color
                    if fill_boxes:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = fill_color
                else:
                    # Remove border
                    shape.line.width = Pt(0)
                    shape.line.fill.background()

                    # Hide text if present
                    if shape.has_text_frame:
                        shape.text_frame.text = ""
            except Exception as e:
                print(f"Could not process shape '{shape_name}' (type={shape.shape_type}): {e}")

    return labeled_elements
